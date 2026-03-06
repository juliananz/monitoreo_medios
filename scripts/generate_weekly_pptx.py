"""
Weekly PPTX report generator for media monitoring.
Secretaría de Economía — Gobierno de Coahuila

Usage:
    python scripts/generate_weekly_pptx.py
    python scripts/generate_weekly_pptx.py --fecha-fin 2026-02-02

Slides:
    1. Portada
    2. Resumen de la Semana (KPIs + topic bars)
    3-N. One slide per topic (headlines)
    N+1. Coahuila: Actividad Regional esta Semana  ← new
    N+2. Conclusiones
"""

import sys
import unicodedata
from collections import defaultdict
from datetime import date, timedelta
from pathlib import Path
from typing import Optional

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))
from config.settings import GROQ_API_KEY, GROQ_MODEL, OUTPUT_DIR

DASHBOARD_CSV = Path(OUTPUT_DIR) / "dashboard_noticias.csv"
XLSX_PATH = Path(__file__).resolve().parents[1] / "regiones_coahuila.xlsx"


# =============================================================================
# STEP 2: detect_coahuila_region
# =============================================================================

def _norm(text: str) -> str:
    """Lowercase and strip accents."""
    text = str(text).lower()
    text = unicodedata.normalize("NFD", text)
    return "".join(c for c in text if unicodedata.category(c) != "Mn")


def _cargar_regiones_xlsx() -> tuple[dict, list]:
    """
    Load region -> [municipality keywords] mapping from regiones_coahuila.xlsx.

    Each municipality name is normalized (lowercase, no accents) and added as
    a keyword. Compound names with punctuation (e.g. 'Francisco I. Madero')
    also get a period-stripped variant ('francisco i madero').

    Returns:
        (regiones_dict, region_order_list)
    """
    df = pd.read_excel(XLSX_PATH)
    df.columns = [c.strip() for c in df.columns]

    regiones: dict = {}
    order: list = []

    for _, row in df.iterrows():
        municipio = str(row["Municipio"]).strip()
        region = str(row["Región"]).strip()
        if not municipio or municipio == "nan" or not region or region == "nan":
            continue

        kw = _norm(municipio)
        kw_alt = kw.replace(".", "").replace("  ", " ").strip()

        if region not in regiones:
            regiones[region] = []
            order.append(region)

        regiones[region].append(kw)
        if kw_alt != kw:
            regiones[region].append(kw_alt)

    return regiones, order


# Load at module level — raises clearly if xlsx is missing or malformed
REGIONES_COAHUILA, REGION_ORDER = _cargar_regiones_xlsx()

# Secondary keyword layer for Carbonífera: fires only when no municipality
# matched AND nivel_geografico is "estatal" or "nacional".
CARBONIFERA_KEYWORDS = [
    "carbon",
    "carbonera",
    "minera",
    "mina",
    "ahmsa",
    "altos hornos",
    "hulla",
    "coque",
    "minero",
    "extraccion",
    "yacimiento",
    "sabinas basin",
    "cuenca carbonifera",
]

_CARBONIFERA_GEO_LEVELS = {"estatal", "nacional"}


def detect_coahuila_region(titulo: str, lugares: str,
                           nivel_geografico: str = "") -> Optional[str]:
    """
    Detect Coahuila sub-region from article fields.

    Priority 1 — Municipality match (all regions):
        Normalizes titulo + lugares and checks each region's municipality
        keywords as substrings. Returns the first match found.

    Priority 2 — Carbonífera topic keywords (secondary layer):
        Only attempted when:
          - No municipality matched in Priority 1, AND
          - nivel_geografico is "estatal" or "nacional"
        Checks CARBONIFERA_KEYWORDS against the normalized titulo.
        Returns "Carbonífera" on the first keyword hit.

    Args:
        titulo:           article headline
        lugares:          NER-extracted locations (may be empty/NaN)
        nivel_geografico: geographic scope label from the pipeline

    Returns:
        Region name string or None
    """
    lugares_clean = lugares if pd.notna(lugares) and str(lugares) not in ("nan", "") else ""
    text = _norm(titulo) + " " + _norm(lugares_clean)

    # Priority 1: municipality name match
    for region, keywords in REGIONES_COAHUILA.items():
        for kw in keywords:
            if kw in text:
                return region

    # Priority 2: Carbonífera topic keywords (geo-gated)
    if _norm(nivel_geografico) in _CARBONIFERA_GEO_LEVELS:
        titulo_norm = _norm(titulo)
        for kw in CARBONIFERA_KEYWORDS:
            if kw in titulo_norm:
                return "Carbonífera"

    return None


# =============================================================================
# DATA LOADING
# =============================================================================

def _cargar_csv() -> pd.DataFrame:
    df = pd.read_csv(DASHBOARD_CSV)
    df["fecha"] = pd.to_datetime(df["fecha"], utc=True, errors="coerce").dt.tz_localize(None)
    df = df.fillna("")
    df["riesgo"] = pd.to_numeric(df["riesgo"], errors="coerce").fillna(0).astype(int)
    df["oportunidad"] = pd.to_numeric(df["oportunidad"], errors="coerce").fillna(0).astype(int)
    return df


def _semana_df(df: pd.DataFrame, fecha_fin: date) -> pd.DataFrame:
    """Return rows for the 7-day window ending on fecha_fin."""
    fecha_ini = pd.Timestamp(fecha_fin - timedelta(days=6))
    fecha_fin_ts = pd.Timestamp(fecha_fin)
    return df[(df["fecha"] >= fecha_ini) & (df["fecha"] <= fecha_fin_ts)].copy()


# =============================================================================
# STEP 3: Regional stats
# =============================================================================

def calcular_stats_regionales(df_week: pd.DataFrame) -> dict:
    """
    Apply detect_coahuila_region to every article in df_week.

    Returns dict keyed by region name:
        {count, articulos: [{titulo, riesgo, oportunidad}], riesgos, oportunidades}
    """
    stats = {
        r: {"count": 0, "articulos": [], "riesgos": 0, "oportunidades": 0}
        for r in REGION_ORDER
    }

    for _, row in df_week.iterrows():
        region = detect_coahuila_region(
            row["titulo"],
            row.get("lugares", ""),
            row.get("nivel_geografico", ""),
        )
        if region:
            stats[region]["count"] += 1
            stats[region]["articulos"].append({
                "titulo": row["titulo"],
                "riesgo": int(row["riesgo"]),
                "oportunidad": int(row["oportunidad"]),
            })
            stats[region]["riesgos"] += int(row["riesgo"])
            stats[region]["oportunidades"] += int(row["oportunidad"])

    return stats


# =============================================================================
# GROQ REGIONAL OUTLOOK
# =============================================================================

def generar_outlook_regional(stats: dict, fecha_fin: date) -> str:
    """
    Call GROQ to generate a 2-sentence regional outlook:
    '¿Qué región concentró más actividad y qué implica para inversión en Coahuila?'
    Falls back to a static summary if GROQ is unavailable.
    """
    resumen_datos = ", ".join(
        f"{r}: {stats[r]['count']} artículos"
        for r in REGION_ORDER
        if stats[r]["count"] > 0
    ) or "Sin menciones regionales esta semana."

    if not GROQ_API_KEY:
        top = max(stats, key=lambda r: stats[r]["count"])
        return (
            f"La región {top} concentró la mayor actividad noticiosa esta semana "
            f"con {stats[top]['count']} menciones. "
            "(Análisis GROQ no disponible — GROQ_API_KEY no configurado.)"
        )

    carbonifera_count = stats.get("Carbonífera", {}).get("count", 0)
    carbonifera_note = (
        " Si la región Carbonífera tiene menciones, enmarca su actividad "
        "en términos de empleo, transición energética y reconversión industrial "
        "(no como inversión general). Para el resto de las regiones, "
        "enmarca la actividad en términos de inversión e industria."
        if carbonifera_count > 0
        else ""
    )

    prompt = (
        "Eres analista de la Secretaría de Economía de Coahuila. "
        f"La semana que termina el {fecha_fin.isoformat()} registró las siguientes "
        f"menciones regionales en medios: {resumen_datos}. "
        "Responde en exactamente 2 oraciones, en español institucional, sin introducción: "
        "¿Qué región concentró más actividad noticiosa y qué implica para "
        f"Coahuila esta semana?{carbonifera_note}"
    )

    try:
        from groq import Groq
        client = Groq(api_key=GROQ_API_KEY)
        response = client.chat.completions.create(
            model=GROQ_MODEL,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.3,
            max_tokens=180,
        )
        return response.choices[0].message.content.strip()
    except Exception as exc:
        return f"Perspectiva regional no disponible ({exc})."


# =============================================================================
# PPTX DESIGN CONSTANTS
# =============================================================================

AZUL   = RGBColor(0x1F, 0x38, 0x64)
TEAL   = RGBColor(0x17, 0x65, 0x8A)
VERDE  = RGBColor(0x1E, 0x7A, 0x45)
ROJO   = RGBColor(0xC0, 0x00, 0x00)
GRIS   = RGBColor(0x99, 0x99, 0x99)
NEGRO  = RGBColor(0x1A, 0x1A, 0x1A)
BLANCO = RGBColor(0xFF, 0xFF, 0xFF)
AZUL_CLARO = RGBColor(0xF0, 0xF4, 0xF8)
GRIS_CLARO = RGBColor(0xE8, 0xE8, 0xE8)

W = Inches(10)     # 16:9 widescreen
H = Inches(5.63)


# =============================================================================
# PPTX LOW-LEVEL HELPERS
# =============================================================================

def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])  # Blank


def _rect(slide, x, y, w, h, fill: Optional[RGBColor] = None):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # 1 = Rectangle
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    return shape


def _tb(slide, text, x, y, w, h,
        bold=False, size=12, color: RGBColor = NEGRO,
        align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold = bold
    run.font.size = Pt(size)
    run.font.color.rgb = color
    return txb


def _header(slide, title: str, subtitle: str = ""):
    _rect(slide, Inches(0), Inches(0), W, Inches(0.72), fill=AZUL)
    _tb(slide, title,
        Inches(0.2), Inches(0.06), Inches(7.2), Inches(0.62),
        bold=True, size=17, color=BLANCO)
    if subtitle:
        _tb(slide, subtitle,
            Inches(7.2), Inches(0.15), Inches(2.6), Inches(0.45),
            size=9, color=RGBColor(0xCC, 0xCC, 0xCC), align=PP_ALIGN.RIGHT)


def _footer(slide, fecha_str: str):
    _rect(slide, Inches(0), Inches(5.4), W, Inches(0.23), fill=GRIS_CLARO)
    _tb(slide,
        f"Secretaría de Economía — Gobierno de Coahuila  |  {fecha_str}",
        Inches(0.2), Inches(5.4), Inches(9.6), Inches(0.23),
        size=7, color=GRIS)


# =============================================================================
# SLIDE BUILDERS
# =============================================================================

def _add_portada(prs: Presentation, fecha_ini: date, fecha_fin: date):
    slide = _blank(prs)
    _rect(slide, Inches(0), Inches(0), W, H, fill=AZUL)
    _tb(slide, "Reporte Semanal\nMonitoreo de Medios",
        Inches(1), Inches(1.0), Inches(8), Inches(2.0),
        bold=True, size=32, color=BLANCO, align=PP_ALIGN.CENTER)
    _tb(slide, "Secretaría de Economía — Gobierno de Coahuila",
        Inches(1), Inches(2.9), Inches(8), Inches(0.55),
        size=14, color=RGBColor(0xAA, 0xCC, 0xEE), align=PP_ALIGN.CENTER)
    _tb(slide,
        f"{fecha_ini.strftime('%d %b')} – {fecha_fin.strftime('%d %b %Y')}",
        Inches(1), Inches(3.48), Inches(8), Inches(0.45),
        size=12, color=RGBColor(0xCC, 0xDD, 0xEE), align=PP_ALIGN.CENTER)


def _add_resumen_slide(prs: Presentation, df_week: pd.DataFrame, fecha_fin: date):
    slide = _blank(prs)
    _header(slide, "Resumen de la Semana", fecha_fin.strftime("%d/%m/%Y"))
    _footer(slide, fecha_fin.strftime("%d/%m/%Y"))

    total  = len(df_week)
    riesgos = int(df_week["riesgo"].sum())
    oports  = int(df_week["oportunidad"].sum())
    medios  = df_week["medio"].nunique()

    kpis = [
        ("Noticias relevantes", str(total),   AZUL),
        ("Riesgos detectados",  str(riesgos), ROJO),
        ("Oportunidades",       str(oports),  VERDE),
        ("Medios activos",      str(medios),  TEAL),
    ]
    for i, (label, value, color) in enumerate(kpis):
        x = Inches(0.4 + i * 2.35)
        _rect(slide, x, Inches(0.9), Inches(2.1), Inches(1.35), fill=color)
        _tb(slide, value, x, Inches(0.9), Inches(2.1), Inches(0.88),
            bold=True, size=36, color=BLANCO, align=PP_ALIGN.CENTER)
        _tb(slide, label, x, Inches(1.72), Inches(2.1), Inches(0.38),
            size=9, color=BLANCO, align=PP_ALIGN.CENTER)

    # Topic frequency bars
    _tb(slide, "Temas principales esta semana",
        Inches(0.4), Inches(2.55), Inches(9.2), Inches(0.28),
        bold=True, size=11, color=AZUL)

    temas_raw = df_week[df_week["temas"].astype(str).str.strip() != ""]["temas"].astype(str)
    tema_counts: dict = defaultdict(int)
    for t in temas_raw:
        for tema in t.split("|"):
            tema = tema.strip()
            if tema:
                tema_counts[tema] += 1

    top_temas = sorted(tema_counts.items(), key=lambda x: -x[1])[:8]
    for idx, (tema, cnt) in enumerate(top_temas):
        col = idx % 2
        row = idx // 2
        x = Inches(0.4 + col * 4.8)
        y = Inches(2.9 + row * 0.44)
        pct = cnt / total if total else 0
        _rect(slide, x, y + Inches(0.06), Inches(4.2), Inches(0.27),
              fill=RGBColor(0xE0, 0xEA, 0xF4))
        bar_w = max(Inches(0.02), Inches(4.2 * pct))
        _rect(slide, x, y + Inches(0.06), bar_w, Inches(0.27), fill=TEAL)
        _tb(slide, f"{tema.replace('_', ' ').title()}  {cnt}",
            x + Inches(0.05), y, Inches(4.1), Inches(0.36), size=9, color=NEGRO)


def _add_topic_slide(prs: Presentation, tema: str,
                     df_tema: pd.DataFrame, fecha_fin: date):
    slide = _blank(prs)
    n = len(df_tema)
    _header(slide, f"Tema: {tema.replace('_', ' ').title()}",
            f"{n} artículo{'s' if n != 1 else ''} esta semana")
    _footer(slide, fecha_fin.strftime("%d/%m/%Y"))

    riesgos = int(df_tema["riesgo"].sum())
    oports  = int(df_tema["oportunidad"].sum())
    _tb(slide, f"{oports} oportunidades  |  {riesgos} riesgos",
        Inches(0.4), Inches(0.78), Inches(9.2), Inches(0.28),
        size=10, color=TEAL)

    top_rows = df_tema.sort_values("fecha", ascending=False).head(8)
    y = Inches(1.18)
    for _, row in top_rows.iterrows():
        titulo = str(row["titulo"])[:110]
        medio  = str(row.get("medio", ""))
        try:
            fecha_str = row["fecha"].strftime("%d/%m") if pd.notna(row["fecha"]) else ""
        except Exception:
            fecha_str = ""

        if row["riesgo"]:
            dot_color = ROJO
        elif row["oportunidad"]:
            dot_color = VERDE
        else:
            dot_color = NEGRO

        _tb(slide, f"● {titulo}", Inches(0.4), y, Inches(8.1), Inches(0.32),
            size=9, color=dot_color)
        _tb(slide, f"{medio}  {fecha_str}",
            Inches(8.5), y, Inches(1.3), Inches(0.32),
            size=7, color=GRIS, align=PP_ALIGN.RIGHT)
        y += Inches(0.36)


# =============================================================================
# STEP 4: Coahuila regional slide
# =============================================================================

def _add_regional_slide(prs: Presentation, stats: dict, outlook: str,
                        total_matched: int, fecha_fin: date):
    """
    Coahuila sub-regional activity slide — inserted before Conclusiones.

    Layout:
      - Blue header bar
      - Summary line: "X artículos con referencia regional directa esta semana"
      - Left column (0.4"–6.3"): one block per region
          Active regions:  name (bold) + count, risk/opp balance, top 2 titles
          Zero regions:    muted gray line "Sin menciones esta semana"
      - Right panel (6.5"–9.8"): light-blue box with GROQ regional outlook
      - Footer
    """
    slide = _blank(prs)
    _header(slide, "Coahuila: Actividad Regional esta Semana",
            fecha_fin.strftime("%d/%m/%Y"))
    _footer(slide, fecha_fin.strftime("%d/%m/%Y"))

    # Summary line
    _tb(slide,
        f"{total_matched} artículos con referencia regional directa esta semana",
        Inches(0.4), Inches(0.78), Inches(9.2), Inches(0.3),
        bold=True, size=11, color=AZUL)

    # Right panel: GROQ outlook
    # Panel spans from y=1.15 to y=5.25 (just above footer at 5.4)
    panel_x = Inches(6.5)
    _rect(slide, panel_x, Inches(1.15), Inches(3.25), Inches(4.1), fill=AZUL_CLARO)
    _tb(slide, "Perspectiva Regional",
        panel_x + Inches(0.15), Inches(1.22), Inches(2.9), Inches(0.30),
        bold=True, size=10, color=AZUL)
    _tb(slide, outlook,
        panel_x + Inches(0.15), Inches(1.55), Inches(2.92), Inches(3.6),
        size=9, color=NEGRO, wrap=True)

    # Left column: region blocks
    # Budget: usable height = footer (5.4") - start (1.18") = 4.22"
    # Worst case: 5 active regions × ~0.82" each = 4.10" — fits.
    # Per active region: name(0.25) + balance(0.17) + title×2(0.20×2) = 0.82"
    # Per inactive region: 0.25"
    y = Inches(1.18)
    for region in REGION_ORDER:
        s = stats[region]

        if s["count"] > 0:
            # Region name + count (two runs in one textbox)
            txb = slide.shapes.add_textbox(Inches(0.4), y, Inches(5.8), Inches(0.28))
            tf = txb.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            r1 = p.add_run()
            r1.text = f"{region}  "
            r1.font.bold = True
            r1.font.size = Pt(10)
            r1.font.color.rgb = AZUL
            r2 = p.add_run()
            r2.text = f"({s['count']} artículo{'s' if s['count'] != 1 else ''})"
            r2.font.bold = False
            r2.font.size = Pt(9)
            r2.font.color.rgb = TEAL
            y += Inches(0.25)

            # Risk / opportunity balance
            bal_parts = []
            if s["oportunidades"]:
                bal_parts.append(
                    f"{s['oportunidades']} oportunidad{'es' if s['oportunidades'] != 1 else ''}"
                )
            if s["riesgos"]:
                bal_parts.append(
                    f"{s['riesgos']} riesgo{'s' if s['riesgos'] != 1 else ''}"
                )
            if bal_parts:
                bal_color = ROJO if s["riesgos"] else VERDE
                _tb(slide, "  " + "  |  ".join(bal_parts),
                    Inches(0.4), y, Inches(5.8), Inches(0.20),
                    size=7, color=bal_color)
                y += Inches(0.17)

            # Top 2 titles
            for art in s["articulos"][:2]:
                titulo_trunc = str(art["titulo"])[:90]
                _tb(slide, f"  • {titulo_trunc}",
                    Inches(0.4), y, Inches(5.9), Inches(0.22),
                    size=7.5, color=NEGRO)
                y += Inches(0.20)

        else:
            # Zero articles — muted line
            _tb(slide,
                f"{region}  —  Sin menciones esta semana",
                Inches(0.4), y, Inches(5.8), Inches(0.24),
                size=8.5, color=GRIS)
            y += Inches(0.25)


def _add_conclusiones_slide(prs: Presentation, df_week: pd.DataFrame, fecha_fin: date):
    slide = _blank(prs)
    _header(slide, "Conclusiones y Próximos Pasos", fecha_fin.strftime("%d/%m/%Y"))
    _footer(slide, fecha_fin.strftime("%d/%m/%Y"))

    total   = len(df_week)
    riesgos = int(df_week["riesgo"].sum())
    oports  = int(df_week["oportunidad"].sum())

    tema_counts: dict = defaultdict(int)
    for t in df_week[df_week["temas"].astype(str).str.strip() != ""]["temas"].astype(str):
        for tema in t.split("|"):
            tema = tema.strip()
            if tema:
                tema_counts[tema] += 1
    top_tema = max(tema_counts, key=tema_counts.get) if tema_counts else "N/D"

    lines = [
        f"• Se analizaron {total} noticias relevantes en la semana.",
        f"• Tema dominante: {top_tema.replace('_', ' ').title()} "
        f"({tema_counts.get(top_tema, 0)} artículos).",
        f"• Balance global: {oports} oportunidades / {riesgos} riesgos detectados.",
        "• Revisar artículos con análisis profundo requerido para seguimiento puntual.",
        "• Próxima actualización: lunes siguiente.",
    ]
    y = Inches(1.15)
    for line in lines:
        _tb(slide, line, Inches(0.6), y, Inches(9.0), Inches(0.42),
            size=13, color=NEGRO)
        y += Inches(0.52)


# =============================================================================
# MAIN GENERATOR
# =============================================================================

def generar_pptx(fecha_fin: Optional[date] = None) -> Path:
    """
    Generate the weekly media monitoring PPTX.

    Args:
        fecha_fin: last day of the reporting week.
                   Defaults to the most recent date with data in the CSV.

    Returns:
        Path to the saved .pptx file.
    """
    df_all = _cargar_csv()

    if fecha_fin is None:
        max_ts = df_all["fecha"].max()
        fecha_fin = max_ts.date() if pd.notna(max_ts) else date.today()

    fecha_ini = fecha_fin - timedelta(days=6)
    df_week = _semana_df(df_all, fecha_fin)

    print(f"Week: {fecha_ini} -> {fecha_fin}  |  {len(df_week)} articles")

    # ------------------------------------------------------------------
    # Step 6: Regional detection preview (printed before generating PPTX)
    # ------------------------------------------------------------------
    stats = calcular_stats_regionales(df_week)
    total_matched = sum(s["count"] for s in stats.values())

    print(f"\nRegional detection results ({total_matched} matched / {len(df_week)} total):")
    for region in REGION_ORDER:
        s = stats[region]
        if s["count"]:
            print(f"  [OK] {region}: {s['count']} articles  "
                  f"({s['oportunidades']} oportunidades, {s['riesgos']} riesgos)")
        else:
            print(f"  [ ] {region}: 0 articles")

    # GROQ regional outlook
    print("\nGenerating GROQ regional outlook...")
    outlook = generar_outlook_regional(stats, fecha_fin)
    print(f"  {outlook[:120]}{'...' if len(outlook) > 120 else ''}")

    # ------------------------------------------------------------------
    # Build PPTX
    # ------------------------------------------------------------------
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    _add_portada(prs, fecha_ini, fecha_fin)
    _add_resumen_slide(prs, df_week, fecha_fin)

    # Topic slides (one per topic, sorted by frequency descending)
    temas_raw = df_week[df_week["temas"].astype(str).str.strip() != ""].copy()
    temas_raw["_tlist"] = temas_raw["temas"].astype(str).str.split("|")
    temas_exp = temas_raw.explode("_tlist")
    temas_exp["_tema"] = temas_exp["_tlist"].str.strip()
    temas_exp = temas_exp[temas_exp["_tema"] != ""]

    tema_freq = temas_exp.groupby("_tema").size().sort_values(ascending=False)
    for tema in tema_freq.index:
        df_tema = temas_exp[temas_exp["_tema"] == tema].drop_duplicates("id")
        _add_topic_slide(prs, tema, df_tema, fecha_fin)

    # ---- Coahuila regional slide (BEFORE Conclusiones) ----
    _add_regional_slide(prs, stats, outlook, total_matched, fecha_fin)

    # ---- Conclusiones ----
    _add_conclusiones_slide(prs, df_week, fecha_fin)

    # Save
    out_dir = Path(OUTPUT_DIR)
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / f"reporte_semanal_{fecha_fin.isoformat()}.pptx"
    prs.save(str(out_path))
    print(f"\nPPTX saved: {out_path}  ({prs.slides.__len__()} slides)")
    return out_path


# =============================================================================
# CLI
# =============================================================================

if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(
        description="Generate weekly media monitoring PPTX — SEC Coahuila"
    )
    parser.add_argument(
        "--fecha-fin", type=str, default=None,
        help="Last date of the report week (YYYY-MM-DD). "
             "Defaults to most recent date with data in the CSV.",
    )
    args = parser.parse_args()

    target_date = None
    if args.fecha_fin:
        try:
            target_date = date.fromisoformat(args.fecha_fin)
        except ValueError:
            print(f"Invalid date: {args.fecha_fin}. Use YYYY-MM-DD.")
            sys.exit(1)

    path = generar_pptx(target_date)
    print(f"\nDone: {path}")
